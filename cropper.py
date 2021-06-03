import os
import fitz
import numpy as np
import argparse
import traceback
import io
from PIL import Image


def clip(value, min_value, max_value):
    if min_value is not None:
        if value < min_value:
            value = min_value
    if max_value is not None:
        if value > max_value:
            value = max_value
    return value

def crop(args):
    infile = args.input
    # process filepath
    infile = os.path.join(os.getcwd(), infile)
    infile = '\\'.join(infile.split('/'))
    inbase = infile.rsplit('.', maxsplit=1)[0]

    outfile = args.output
    if outfile is None:
        outfile = inbase + '_crop.pdf'
    outbase = outfile.rsplit('.', maxsplit=1)[0]
    names = None
    thresh = args.thresh

    # parse working dir.
    workdir = '/'.join('/'.join(infile.split('\\')).split('/')[:-1])
    if workdir == '':
        workdir = '.'

    if infile.split('.')[-1].lower() in ['ppt', 'pptx']:
        # lets process pptx.
        import comtypes.client
        from pptx import Presentation
        # convert to pdf.
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        slides = powerpoint.Presentations.Open(infile)
        print('convert {0} to {1}'.format(infile, inbase + '.pdf'))
        slides.SaveAs(inbase + '.pdf', 32)
        slides.close()
        # powerpoint.Quit()
        # read comments.
        ppt = Presentation(infile)
        slides = ppt.slides
        names = []
        for i, slide in enumerate(slides):
            anno = slide.notes_slide.notes_text_frame.text.strip()
            if anno == '':
                anno = '{0}_{1}.pdf'.format(outbase, i+1)
            names.append(anno)
        names = [n + '.pdf' if n[:-4] != '.pdf' else n for n in names]
        infile = inbase + '.pdf'

    if infile.split('.')[-1].lower() != 'pdf':
        infile += '.pdf'


    pdf = fitz.open(infile)
    for i, page in enumerate(pdf):
        try:
            # recover crop box
            page.setCropBox(page.MediaBox)
            # render page as img
            zoom = args.zoom
            mat = fitz.Matrix(zoom, zoom)
            img = page.getPixmap(matrix=mat)

            # convert the img to nparray.
            # img = [[img.pixel(j, i) for j in range(img.w)] for i in range(img.h)]
            img = Image.open(io.BytesIO(img.tobytes()))
            img = np.asarray(img, dtype=np.float32)

            # get ltrb
            background = np.broadcast_to(args.background_color, (img.shape))
            diff = np.abs(img - background).mean(axis=2)
            points = np.argwhere(diff > thresh)
            x = points[:, 1]
            y = points[:, 0]
            left, right = min(x), max(x)
            top, bottom = min(y), max(y)
            right += 1
            bottom += 1

            # calculate real crop coordinate.
            ltrb = left, top, right, bottom
            ltrb = [coord/zoom for coord in ltrb]

            # reserve border.
            border = args.border
            left, top, right, bottom = [x+y for x, y in zip(ltrb, [-border,-border,border,border])]

            # clip ltrb
            left = clip(left, 0, img.shape[1]/zoom)
            top = clip(top, 0, img.shape[0]/zoom)
            right = clip(right, 0, img.shape[1]/zoom)
            bottom = clip(bottom, 0, img.shape[0]/zoom)
            if args.visual:
                print('cropbox =', (left, top), (right, bottom))
            page.setCropBox(fitz.Rect(left, top, right, bottom))
        except KeyboardInterrupt:
            print('Operation Cancled.')
            break
        except:
            print('Error happened while cropint page', i+1)
            print(traceback.format_exc())

    if not args.split:
        pdf.save(outfile)
    else:
        if names is None:
            names = ['{0}_{1}.pdf'.format(outbase, i+1) for i in range(len(pdf))]
        if args.names is not None:
            names = args.names
            namesfile = os.path.join(workdir, names)
            if os.path.isfile(namesfile) and names.split('.')[-1] == 'txt':
                with open(namesfile, 'r', encoding='utf-8') as f:
                    names = f.read()
                names = names.split('\n')
                names = [n.strip() for n in names]
                names = [n for n in names if n != '']
            else:
                names = names.split(',')
            names = [n + '.pdf' if n[:-4] != '.pdf' else n for n in names]
        if len(names) < len(pdf):
            names += ['{0}_{1}.pdf'.format(outbase, i+1) for i in range(len(names), len(pdf))]
        names = [os.path.join(workdir, name) for name in names]
        for i, (name, page) in enumerate(zip(names, pdf)):
            page_pdf = fitz.Document()
            page_pdf.insertPDF(pdf, from_page=i, to_page=i)
            page_pdf.save(name)
            if not args.mute:
                print('saving cropped page to', name)
    pdf.close()